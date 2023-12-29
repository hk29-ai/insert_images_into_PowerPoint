import glob
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# 画像ファイルのパスのリスト
image_paths = glob.glob("./dir_img/*.jpg")
num_figs = len(image_paths)
print(num_figs, image_paths)

# スライド1枚に画像を貼り付ける列数
num_cols = 2

# 保存するファイル名
save_file_name = f"{num_cols}_cols.pptx"

# 画像サイズを取得する関数
def get_image_size(image_path):
    # 画像サイズを取得
    with Image.open(image_path) as img:
        width, height = img.size

    return width, height

# 画像サイズを取得
image_width, image_height = get_image_size(image_paths[0])
# オリジナル画像のアスペクト比（縦横比）を取得
original_aspect_ratio = image_height / image_width

# 新規プレゼンテーションを作成
presentation = Presentation()

# スライドの幅と高さを取得
slide_width = presentation.slide_width - (2 * Inches(0.5))  # 左右の余白を引く
slide_height = presentation.slide_height - (2 * Inches(0.5))  # 上下の余白を引く

# 行列の間隔と行の上部のスペース、左右のスペースを設定
row_spacing = Inches(0.3)  # 行間
col_spacing = Inches(0.3)  # 列間
top_space_first_row = Inches(1.5)  # 上部スペースを持たせる
side_space = Inches(1)  # 左スペースを持たせる

# 行数の計算
if (len(image_paths) % num_cols) == 0:
    num_rows = len(image_paths) // num_cols
else:
    num_rows = (len(image_paths) // num_cols) + 1

# 画像を横に並べるための計算
total_side_space = 2 * side_space  # 左右のスペースの和
#image_and_spacing_width = num_cols * (slide_width - total_side_space) + (num_cols - 1) * col_spacing
image_and_spacing_width = num_cols * (slide_width - total_side_space) \
                          + (num_cols - 1) * col_spacing

# 画像のサイズを調整
if image_and_spacing_width > slide_width:
    # スライドの幅に収まらない場合、画像サイズを調整
    total_col_spacing = (num_cols - 1) * col_spacing
    adjusted_image_width = (slide_width - total_side_space - total_col_spacing) / num_cols
    adjusted_image_height =  adjusted_image_width * original_aspect_ratio
    
    # 縦方向にはみ出さないための行数を算出
    for i in range(1, num_rows+1, 1):
        image_and_spacing_height = adjusted_image_height * i + top_space_first_row + row_spacing*(i-1)
        height_delta = slide_height - image_and_spacing_height
        if height_delta < 0:
            num_rows = i
            break
elif num_cols == 1: # 1ページに1枚の場合
    total_col_spacing = (num_cols - 1) * col_spacing
    adjusted_image_width = (slide_width - total_side_space - total_col_spacing) / num_cols
    adjusted_image_height =  adjusted_image_width * original_aspect_ratio
    num_rows = 1

print('num_cols, num_rows', num_cols, num_rows)    

# 必要なスライドページ数を計算
if (num_figs % (num_rows * num_cols)) == 0:
    pages = num_figs // (num_rows * num_cols)            
else:
    pages = num_figs // (num_rows * num_cols)  + 1  
    
# スライドを必要な分だけ作成する
slide_count = 0
fig_count = 0
while slide_count < pages:
    # スライドを追加
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "タイトル"
    # 左詰め
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # タイトルのフォントサイズを32ポイントに設定
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    # タイトルをスライドの上端に配置
    title_shape.top = Inches(0.1)
    # タイトルの幅をスライドの幅に合わせる
    title_shape.width = presentation.slide_width
    # タイトルの高さをフォントサイズと同じに設定
    title_shape.height = title_shape.text_frame.paragraphs[0].font.size
    
    # サブタイトルを設定
    left = title_shape.left
    top = title_shape.top + title_shape.height
    width = title_shape.width
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.text = "サブタイトル"
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 中央配置
    frame.paragraphs[0].runs[0].font.size = Pt(28)  # フォントサイズを24に設定
    
    # タイトルの下に棒線を描く
    left = title_shape.left
    top = title_shape.top + title_shape.height
    width = title_shape.width
    height = Pt(1)  # 線の太さを1ptに設定

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    fill = line.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)  # 黒
    line.shadow.inherit = True  # 影
    
    # 画像ファイルを貼り付けてゆく
    if num_rows != 1: # スライド1枚に複数の画像を貼り付ける場合
        for row in range(num_rows):
            for col in range(num_cols):
                #print('row, col', row, col)
                idx = row * num_cols + col
                #print(idx)
                left_inch = col * (adjusted_image_width + col_spacing) + side_space

                top_inch = top_space_first_row \
                           + row * (adjusted_image_width * original_aspect_ratio + row_spacing)

                # 画像を追加
                fig_count = idx + slide_count*(num_rows*num_cols)
                picture = slide.shapes.add_picture(image_paths[fig_count],
                                                   left_inch,
                                                   top_inch,
                                                   width=adjusted_image_width)

                # 画像の縦横比を保つ
                picture.lock_aspect_ratio = True

                # １ページあたりの必要枚数を貼り付けたらループを抜ける
                if (idx == (num_rows*num_cols)-1) or (fig_count == num_figs-1):
                    break
            # ２重ループを抜ける場合は次の３行で出来る
            else: 
                continue
            break # ひとつ目のループをbreakで抜けた場合にここにきて、breakにより２つ目のループを抜ける
    else: # スライド1枚に1枚の画像を貼り付ける場合
        left_inch = side_space
        top_inch = top_space_first_row
        # 画像を追加
        picture = slide.shapes.add_picture(image_paths[slide_count],
                                           left_inch,
                                           top_inch,
                                           width=adjusted_image_width)

        # 画像の縦横比を保つ
        picture.lock_aspect_ratio = True
        
    slide_count += 1
    #print('finished slide',slide_count)
        
# プレゼンテーションを保存
presentation.save(save_file_name)
print(f"save {save_file_name}")
